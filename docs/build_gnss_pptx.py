# -*- coding: utf-8 -*-
"""Genere le PPTX technique GNSS double recepteur (Tersus / CHCNAV).

Source de verite : docs/presentation_gnss_technique_tersus_chcnav.md
Sortie          : docs/SRM_Collecte_GNSS_Technique.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette (alignee sur le deck NotebookLM)
DARK = RGBColor(0x1E, 0x29, 0x3B)      # bleu nuit titres
GREY = RGBColor(0x5A, 0x66, 0x75)      # gris texte secondaire
TERSUS = RGBColor(0xF1, 0x5A, 0x3B)    # orange Tersus
CHCNAV = RGBColor(0x1C, 0xA9, 0xC9)    # cyan CHCNAV
LIGHT = RGBColor(0xF2, 0xF4, 0xF6)     # fond clair
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x26, 0x2B, 0x33)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))


def set_text(tf, runs, size, color, bold=False, align=PP_ALIGN.LEFT,
             font="Calibri", space_after=6, line=1.0):
    """runs: str ou liste de (texte, color|None, bold|None)."""
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for txt, c, b in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = b if b is not None else bold
        r.font.color.rgb = c if c is not None else color
        r.font.name = font
    return p


def add_para(tf, runs, size, color, bold=False, align=PP_ALIGN.LEFT,
             font="Calibri", space_after=6, line=1.0, bullet_color=None):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for txt, c, b in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = b if b is not None else bold
        r.font.color.rgb = c if c is not None else color
        r.font.name = font
    return p


def title_band(slide, n, total, title):
    """Bandeau titre standard pour slides de contenu."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.55); tf.margin_top = Inches(0.18)
    tf.word_wrap = True
    set_text(tf, title, 26, WHITE, bold=True)
    # pastille numero
    num = box(slide, 12.1, 0.32, 1.0, 0.5)
    set_text(num.text_frame, f"{n}/{total}", 14, RGBColor(0x9A,0xA6,0xB5),
             bold=True, align=PP_ALIGN.RIGHT)


def code_block(slide, l, t, w, h, lines):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CODE_BG
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.16)
    tf.word_wrap = True
    set_text(tf, lines[0], 11, CODE_FG, font="Consolas", space_after=0, line=1.05)
    for ln in lines[1:]:
        add_para(tf, ln, 11, CODE_FG, font="Consolas", space_after=0, line=1.05)


def table(slide, l, t, w, h, data, col_w=None, header_fill=DARK,
          font_size=12, brand_header=False):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(l), Inches(t),
                                Inches(w), Inches(h)).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Inches(cw)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size); r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.bold = True; r.font.color.rgb = WHITE
                if brand_header and ci == 1:
                    cell.fill.fore_color.rgb = TERSUS
                if brand_header and ci == 2:
                    cell.fill.fore_color.rgb = CHCNAV
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = DARK
    return gt


# ----------------------------------------------------------------------------
# SLIDE 1 - Titre
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, LIGHT)
t = box(s, 1.0, 2.2, 11.3, 2.0)
set_text(t.text_frame, "Integration GNSS double recepteur", 40, DARK,
         bold=True, align=PP_ALIGN.CENTER)
add_para(t.text_frame, [("Tersus", TERSUS, True), (" / ", DARK, True),
                        ("CHCNAV", CHCNAV, True), ("  -  dossier technique", DARK, True)],
         40, DARK, bold=True, align=PP_ALIGN.CENTER)
sub = box(s, 1.0, 4.3, 11.3, 1.2)
set_text(sub.text_frame,
         "Pipeline NMEA -> Merchich, reduction antenne par marque, preuves de validation",
         18, GREY, align=PP_ALIGN.CENTER)
add_para(sub.text_frame,
         "Cale sur le code : antenna_catalog.dart, gnss_config_service.dart, projection_service.dart",
         13, GREY, align=PP_ALIGN.CENTER)
# barre couleur
b1 = s.shapes.add_shape(1, Inches(5.2), Inches(3.95), Inches(1.45), Inches(0.10))
b1.fill.solid(); b1.fill.fore_color.rgb = TERSUS; b1.line.fill.background()
b2 = s.shapes.add_shape(1, Inches(6.65), Inches(3.95), Inches(1.45), Inches(0.10))
b2.fill.solid(); b2.fill.fore_color.rgb = CHCNAV; b2.line.fill.background()

TOTAL = 10

# ----------------------------------------------------------------------------
# SLIDE 2 - These technique
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 2, TOTAL, "These technique")
c = box(s, 0.7, 1.5, 11.9, 5.5); tf = c.text_frame; tf.word_wrap = True
set_text(tf, [("L'app reimplemente bit pour bit la chaine de reduction de ", DARK, False),
              ("Nuwa", TERSUS, True), (" (Tersus) et de ", DARK, False),
              ("LandStar", CHCNAV, True), (" (CHCNAV).", DARK, False)], 18, DARK, line=1.1)
add_para(tf, "Le pipeline de positionnement est IDENTIQUE entre les deux marques.",
         18, DARK, bold=True, space_after=14, line=1.1)
add_para(tf, "Seules deux choses divergent, et sont branchees par marque :", 17, DARK,
         space_after=10)
add_para(tf, [("1.  ", TERSUS, True),
              ("Les constantes geometriques d'antenne (catalogue par marque).", DARK, False)],
         17, DARK, space_after=8)
add_para(tf, [("2.  ", TERSUS, True),
              ("Les formules de reduction APC -> sol (offset vertical).", DARK, False)],
         17, DARK, space_after=16)
add_para(tf, [("Tout le reste (NMEA -> WGS84 -> Merchich) est commun et valide empiriquement a ", DARK, False),
              ("0.1 mm", CHCNAV, True), (" contre Nuwa.", DARK, False)], 17, DARK, line=1.1)

# ----------------------------------------------------------------------------
# SLIDE 3 - Pipeline complet
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 3, TOTAL, "Le pipeline complet (commun aux deux marques)")
code_block(s, 0.6, 1.45, 8.5, 5.3, [
    "Trame NMEA GGA (rover, Bluetooth)",
    "   | champ 9  : altitude orthometrique H",
    "   | champ 11 : separation geoidale N",
    "   v",
    "h_ellipsoidale = H + N      <- ellipsoidal, PAS l'ortho",
    "   v",
    "proj4dart : WGS84(lon,lat,h) -> Merchich EPSG:26191",
    "   | Helmert towgs84 = 31,146,47,0,0,0,0",
    "   | WGS84 -> Clarke 1880",
    "   | AUCUN geoide (.ggf) : pur datum (= Nuwa BLH84_XYh)",
    "   v",
    "(X, Y, Z_apc) metres Merchich   <- Z = centre de phase",
    "   v",
    "Z_sol = Z_apc - apcToGroundVerticalOffset   <- BRANCHE MARQUE",
    "   v",
    "VerticalAdjustment.apply()   None | Constante | Plan incline",
    "   v",
    "(X, Y, Z_sol) stocke",
])
side = box(s, 9.3, 1.6, 3.5, 5.0); tf = side.text_frame; tf.word_wrap = True
set_text(tf, "References code", 15, DARK, bold=True, space_after=8)
add_para(tf, "projection_service.dart:41-88", 12, CHCNAV, font="Consolas", space_after=4)
add_para(tf, "merchichPointFromGnss, wgs84HeightToMerchich", 11, GREY, space_after=14)
add_para(tf, "Soustraction antenne AU SAVE (Z stocke = h_sol), comme LandStar.", 13, DARK, space_after=8, line=1.05)
add_para(tf, "Tersus/Nuwa stocke h_apc et soustrait a l'export : on a aligne sur LandStar pour un Z deja exploitable en base.", 13, GREY, line=1.05)

# ----------------------------------------------------------------------------
# SLIDE 4 - Datum (preuve)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 4, TOTAL, "Pourquoi le datum tombe juste (preuve, pas analogie)")
c = box(s, 0.7, 1.4, 11.9, 2.2); tf = c.text_frame; tf.word_wrap = True
set_text(tf, [("Notre projection_constants.dart : ", DARK, False),
              ("+towgs84=31,146,47,0,0,0,0", CHCNAV, True)], 16, DARK, space_after=8)
add_para(tf, [("Ecran Nuwa MERCHICH MAROC Zone 1, onglet Datum Trans (capture 2026-05-31) : ", DARK, False),
              ("Bursa-Wolf  Dx=-31 Dy=-146 Dz=-47  Rx=Ry=Rz=0  Scale=0", TERSUS, True)],
         16, DARK, line=1.1)
table(s, 0.7, 3.7, 11.9, 1.6, [
    ["", "Rotations / Echelle", "Translation", "Signe"],
    ["Verdict", "nulles des 2 cotes -> Bursa 7p degenere en Helmert 3p (pure translation)",
     "identique", "-31 vs +31 = convention de sens (inverse = negation exacte)"],
], col_w=[1.3, 4.6, 2.0, 4.0], font_size=12)
c2 = box(s, 0.7, 5.6, 11.9, 1.5); tf = c2.text_frame; tf.word_wrap = True
set_text(tf, [("Consequence : ", DARK, True),
              ("l'aller-retour RTK se referme au cm, horizontal ET vertical. Confirme empiriquement : proj4dart 3D = Nuwa BLH->NEH a 0.1 mm.", DARK, False)],
         15, DARK, line=1.1)

# ----------------------------------------------------------------------------
# SLIDE 5 - Reduction APC -> sol (formules)
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 5, TOTAL, "Le seul point de divergence : reduction APC -> sol")
c = box(s, 0.7, 1.25, 11.9, 0.6); tf = c.text_frame; tf.word_wrap = True
set_text(tf, [("GnssRoverConfig.apcToGroundVerticalOffset", CHCNAV, True),
              ("  fait un switch(brand). Formules par methode :", DARK, False)], 13, DARK)
table(s, 0.7, 1.9, 11.9, 1.7, [
    ["Methode", "CHCNAV (LandStar)", "Tersus (Nuwa)"],
    ["Vertical", "A = H + DH", "A = H  (rien ajoute)"],
    ["PhaseCenter", "A = H", "A = H + AntCenter"],
    ["Slant", "A = sqrt(H2 - R0^2) - H0 + DH", "A = |sqrt(H2 - R^2) + AntCenter - AntBottomHeight|"],
], col_w=[2.0, 4.7, 5.2], font_size=12.5, brand_header=True)
c2 = box(s, 0.7, 3.85, 5.7, 0.4); set_text(c2.text_frame, "Mapping constantes (memes champs, semantique differente)", 13, DARK, bold=True)
table(s, 0.7, 4.3, 5.9, 1.5, [
    ["Champ", "CHCNAV", "Tersus"],
    ["r0", "R0 (rayon slant)", "AntRadius"],
    ["h0", "H0 (APR->repere)", "AntBottomHeight (0.0)"],
    ["dh", "DH (APR->APC)", "AntCenter"],
], col_w=[1.0, 2.6, 2.3], font_size=11.5)
warn = s.shapes.add_shape(1, Inches(6.9), Inches(4.3), Inches(5.7), Inches(2.4))
warn.fill.solid(); warn.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xE7)
warn.line.color.rgb = TERSUS; warn.line.width = Pt(1.5)
tf = warn.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
set_text(tf, "Le piege metier", 14, TERSUS, bold=True, space_after=6)
add_para(tf, "Meme hauteur saisie (1,80 m) + mauvaise formule = Z decale de 5 a 10 cm (DH vs AntCenter). Invisible, mais faux pour des profondeurs de reseau.", 13, DARK, line=1.05, space_after=8)
add_para(tf, "Garde-fou Slant : si r0<=0 ou H<=r0 (racine imaginaire), on retombe sur Vertical.", 12, GREY, line=1.05)

# ----------------------------------------------------------------------------
# SLIDE 6 - Catalogues
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 6, TOTAL, "Catalogues d'antennes (valeurs reelles, par marque)")
cap = box(s, 0.7, 1.25, 5.5, 0.4)
set_text(cap.text_frame, [("CHCNAV", CHCNAV, True), ("  (CHCNAV/Config/ante.hpc)", GREY, False)], 14, DARK, bold=True)
table(s, 0.7, 1.7, 5.8, 1.7, [
    ["Modele", "r0 (R0)", "h0 (H0)", "dh (DH)"],
    ["GENERIC", "0.0", "0.0", "0.0"],
    ["i89 / Z100", "0.124", "0.07996", "0.0"],
    ["X16 / X16 Pro", "0.124", "0.08133", "0.0"],
], col_w=[2.2, 1.2, 1.3, 1.1], font_size=12, header_fill=CHCNAV)
cap2 = box(s, 6.9, 1.25, 5.7, 0.4)
set_text(cap2.text_frame, [("Tersus", TERSUS, True), ("  (TbAntenna, Nuwa.apk)", GREY, False)], 14, DARK, bold=True)
table(s, 6.9, 1.7, 5.8, 4.4, [
    ["Modele", "r0 (Radius)", "h0 (Bottom)", "dh (Center)"],
    ["GENERIC", "0.0", "0.0", "0.0"],
    ["OSCAR (defaut)", "0.13", "0.0", "0.094"],
    ["AX3702", "0.13", "0.0", "0.054"],
    ["AX3702 (HG)", "0.13", "0.0", "0.0509"],
    ["AX4E02", "0.13", "0.0", "0.059"],
    ["LUKA", "0.13", "0.0", "0.082"],
    ["TS20", "0.13", "0.0", "0.06425"],
    ["K1", "0.13", "0.0", "0.082"],
], col_w=[2.0, 1.3, 1.3, 1.2], font_size=11.5, header_fill=TERSUS)
note = box(s, 0.7, 3.7, 5.8, 2.8); tf = note.text_frame; tf.word_wrap = True
set_text(tf, "Deux maps separees dans antenna_catalog.dart, jamais fusionnees dans l'UI (sortedForBrand filtre par marque).", 14, DARK, line=1.1, space_after=12)
add_para(tf, "Cote Tersus : h0 = 0.0 systematique (catalogue officiel). Tout le decalage APC est porte par AntCenter (dh).", 14, GREY, line=1.1)

# ----------------------------------------------------------------------------
# SLIDE 7 - Architecture logicielle
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 7, TOTAL, "Architecture logicielle de la bascule")
code_block(s, 0.6, 1.45, 7.0, 3.4, [
    "enum GnssBrand { chcnav, tersus }",
    "",
    "GnssRoverConfig {",
    "  brand, antenna, heightMeters,",
    "  surveyType, verticalAdjustment",
    "  apcToGroundVerticalOffset",
    "    -> switch(brand) {",
    "         _chcnavOffset | _tersusOffset }",
    "}",
    "",
    "GnssConfigService -> persistance app_metadata",
])
side = box(s, 7.8, 1.55, 5.0, 5.2); tf = side.text_frame; tf.word_wrap = True
set_text(tf, "Trois defenses", 16, DARK, bold=True, space_after=8)
add_para(tf, [("1. Defaut Tersus assume", TERSUS, True),
              (" : pas de brand persistee = install pre-enum -> Tersus OSCAR (materiel deploye), pas CHCNAV.", DARK, False)], 13, DARK, line=1.05, space_after=10)
add_para(tf, [("2. Anti-mismatch antenne/marque", TERSUS, True),
              (" : antenne hors marque (downgrade, metadata corrompu) -> antenne par defaut de la marque. Jamais la mauvaise formule.", DARK, False)], 13, DARK, line=1.05, space_after=10)
add_para(tf, [("3. Parse tolerant", TERSUS, True),
              (" : VerticalAdjustment au format inconnu retombe sur None au lieu de planter.", DARK, False)], 13, DARK, line=1.05)
foot = box(s, 0.6, 5.1, 7.0, 1.6); tf = foot.text_frame; tf.word_wrap = True
set_text(tf, "Cles app_metadata", 13, DARK, bold=True, space_after=4)
add_para(tf, "gnss_rover_brand / _antenna_model / _antenna_height / _antenna_survey_type / gnss_vertical_adjustment", 11, GREY, font="Consolas", line=1.1)

# ----------------------------------------------------------------------------
# SLIDE 8 - Ajustement vertical
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 8, TOTAL, "Ajustement vertical (au-dessus de la reduction antenne)")
c = box(s, 0.7, 1.3, 11.9, 0.6)
set_text(c.text_frame, [("sealed class VerticalAdjustment", CHCNAV, True),
                        ("  applique APRES la reduction APC->sol, sur l'elevation deja en Merchich :", DARK, False)], 13, DARK)
table(s, 0.7, 2.0, 11.9, 1.8, [
    ["Type", "Formule", "Usage"],
    ["None", "H", "defaut"],
    ["Constante", "H + offset", "calage sur 1 point connu"],
    ["Plan incline", "H + a*X + b*Y + c", "calage multi-points, ecart Z lineaire avec la position"],
], col_w=[2.0, 3.5, 6.4], font_size=13)
c2 = box(s, 0.7, 4.2, 11.9, 2.6); tf = c2.text_frame; tf.word_wrap = True
set_text(tf, "Le plan incline vise un calage moindres carres sur 3+ points connus (outil Phase 3.1).", 15, DARK, line=1.1, space_after=12)
add_para(tf, [("C'est le levier", TERSUS, True),
              (" si le client exige un jour un rattachement borne NGM officielle au cm. La translation pure Nuwa reste une approx ~5 m vs Merchich geodesique officiel, mais la coherence interne et le rattachement au point base sont parfaits tant qu'on ne touche pas a l'onglet Datum Trans.", DARK, False)], 15, DARK, line=1.1)

# ----------------------------------------------------------------------------
# SLIDE 9 - Limites / dette
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 9, TOTAL, "Limites connues et dette assumee")
c = box(s, 0.7, 1.5, 11.9, 5.3); tf = c.text_frame; tf.word_wrap = True
items = [
    ("Pas de geoide (.ggf)", " : volontaire, pour reproduire exactement Nuwa. Z = elevation dans le datum Merchich, pas une altitude orthometrique officielle. Coherent avec le carnet constructeur utilise."),
    ("Approx ~5 m vs Merchich geodesique officiel", " sur l'absolu (translation 3 parametres). Coherence interne et rattachement base = cm."),
    ("L'egalite Nuwa casse", " si : rotation/echelle ajoutees dans Datum Trans, grille NTv2, ou exigence borne NGM. Reponse prevue = calage local (Plan incline) ou grille, pas une refonte du pipeline."),
    ("Doctrine terrain", " : NE PAS modifier l'onglet Datum Trans de Nuwa (garder Bursa -31/-146/-47, rot/scale nuls)."),
]
first = True
for head, body in items:
    runs = [("-  ", TERSUS, True), (head, DARK, True), (body, GREY, False)]
    if first:
        set_text(tf, runs, 16, DARK, line=1.15, space_after=16); first = False
    else:
        add_para(tf, runs, 16, DARK, line=1.15, space_after=16)

# ----------------------------------------------------------------------------
# SLIDE 10 - Synthese
# ----------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title_band(s, 10, TOTAL, "Synthese pour la revue")
table(s, 0.7, 1.5, 11.9, 3.8, [
    ["Dimension", "Etat"],
    ["Pipeline NMEA -> Merchich", "commun, valide 0.1 mm vs Nuwa"],
    ["Datum", "translation pure, identique a Nuwa (preuve ecran + empirique)"],
    ["Reduction antenne", "branchee par marque, formules alignees sur sources decompilees"],
    ["Catalogues", "separes par marque, valeurs verifiees (ante.hpc / TbAntenna)"],
    ["Defenses", "defaut Tersus, anti-mismatch antenne/marque, parse tolerant"],
    ["Calage cm officiel", "non couvert (Plan incline en levier, Phase 3.1)"],
], col_w=[3.8, 8.1], font_size=13)
band = s.shapes.add_shape(1, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.4))
band.fill.solid(); band.fill.fore_color.rgb = DARK; band.line.fill.background()
tf = band.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "Une seule app native qui reimplemente fidelement deux chaines de reduction constructeur : pipeline commun et prouve, divergence (constantes + formule antenne) isolee derriere un enum GnssBrand et un catalogue par marque.",
         15, WHITE, bold=True, align=PP_ALIGN.CENTER, line=1.1)

out = r"c:\Users\AnasDahou\Desktop\srm_collecte\docs\SRM_Collecte_GNSS_Technique.pptx"
prs.save(out)
print("OK ->", out, "|", len(prs.slides._sldIdLst), "slides")
